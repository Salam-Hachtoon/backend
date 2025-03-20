import logging, requests, re
from rest_framework.response import Response # type: ignore
from .models import Attachment
from django.conf import settings
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


# Create a utility logger
logger = logging.getLogger('utility_functions')


def combine_completed_files_content(batch_id):
    """
    Combines the extracted text content of all completed files associated with a given batch ID.

    This function retrieves all attachments linked to the specified batch ID and checks if their
    processing status is marked as "completed". If any file is not completed, an error is logged,
    and a message is returned indicating that not all files have been processed.

    For completed files, the extracted text is concatenated into a single string. If a file has
    an empty extracted text, an error is logged, and the file is skipped.

    Args:
        batch_id (int): The ID of the batch whose attachments are to be processed.

    Returns:
        str: The combined extracted text of all completed files, or a message indicating that
             not all files have been processed.
    """
    for attchemnt in Attachment.objects.filter(batch_id=batch_id):
        if attchemnt.status != "completed":
            logger.error("Not all files have been processed yet.")
            return 'Not all files have been processed yet.'

    combined_text = ""
    for attchemnt in Attachment.objects.filter(batch_id=batch_id):
        if attchemnt.extracted_text == '':
            logger.error("Empty extracted text for file: {}".format(attchemnt.file.name))
            pass
        combined_text += attchemnt.extracted_text + "\n"

    return combined_text



def call_deepseek_ai_summary(combined_text):
    """"
    Generates a structured and well-organized summary of the provided text using the DeepSeek AI API.
    The function sends the input text to the DeepSeek AI API with specific summarization guidelines 
    and retrieves a concise summary in a predefined format.
    Args:
        combined_text (str): The input text to be summarized.
    Returns:
        str: The generated summary if successful, or an error message if the operation fails.
    Raises:
        Exception: Logs and handles any exceptions that occur during the API call.
    Example Output Format:
        "The summary of the input text."
    """
    # Crate the payload for the DeepSeek AI API
    Prompt = """
    You are an advanced AI assistant specialized in text summarization. Your task is to generate a structured and well-organized summary of the user-provided text in **JSON format**.  

    ### **Guidelines:**
    1. Begin with a **clear and relevant title** for the summary.  
    2. Follow it with a **short subtitle** that provides additional context.  
    3. Present the summary as a **single, well-structured paragraph** with all key details.  
    4. Ensure clarity, conciseness, and readability while retaining essential information.  
    5. Avoid unnecessary details and redundant words.  

    ### **Expected JSON Output Format:**
    {
    "summary": {
        "content": "<Summary paragraph with key details in a concise manner>"
    }
    }
    """
    try:
        # Initialize the OpenAI client with the DeepSeek API key and base URL
        client = OpenAI(api_key=settings.DEEPSEEK_API_KEY, base_url="https://api.deepseek.com")

        # Call the DeepSeek API to generate a summary
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": Prompt},
                {"role": "user", "content": combined_text},
            ],
            stream=False
        )

        # Extract the summary from the response
        summary = response.choices[0].message.content
        return summary

    except Exception as e:
        # Log the error and return a failure message
        logger.error("Failed to generate summary: {}".format(e))
        return "Failed to generate summary"


def call_deepseek_ai_flashcards(summary_content):
    """
    Generates a structured set of flashcards in JSON format by analyzing the provided summary content.
    This function sends a request to the DeepSeek AI API with a prompt to generate flashcards based on the 
    given summary. The flashcards include key terms and their definitions, formatted as a JSON object.
    Args:
        summary_content (str): The summary text to be processed for flashcard generation.
    Returns:
        str: A JSON-formatted string containing the generated flashcards if the API call is successful.
             If the API call fails, returns an error message or a default failure message.
    Raises:
        None: The function handles API errors internally and logs them.
    Example:
        summary = "Photosynthesis is the process by which green plants and some other organisms use sunlight to synthesize foods."
        flashcards = call_deepseek_ai_flashcards(summary)
        print(flashcards)
    """

    Prompt = """
    You are an AI assistant skilled in educational content generation. Your task is to analyze the provided summary and generate a structured set of flashcards in JSON format.

    ### **Instructions:**  
    - Extract **key terms** from the summary. These should be important concepts, technical terms, or notable entities.  
    - For each term, provide:  
    - `"term"`: The name of the key concept.  
    - `"definition"`: A concise, clear explanation of the term.  
    - All flashcards must be stored inside a **single key** called `"flashcards"`.  
    - Ensure the output is a **valid JSON object**.

    ### **Text to process:**  
    provided by the user

    ### **Output Format Example:**  
    Return a JSON object structured as follows and remove the markdown symbols from the output:

    ```json
    {
        "flashcards": [
            {
                "term": "Example Term",
                "definition": "A brief explanation of the term."
            },
            {
                "term": "Another Term",
                "definition": "Another brief explanation."
            }
        ]
    }
    """
    try:
        # Initialize the OpenAI client with the DeepSeek API key and base URL
        client = OpenAI(api_key=settings.DEEPSEEK_API_KEY, base_url="https://api.deepseek.com")

        # Call the DeepSeek API to generate a summary
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": Prompt},
                {"role": "user", "content": summary_content},
            ],
            stream=False
        )
        logger.info(response)
        # Extract the summary from the response
        flash_cards = response.choices[0].message.content
        logger.info(flash_cards)
        return flash_cards

    except Exception as e:
        # Log the error and return a failure message
        logger.error("Failed to generate flash cards: {}".format(e))
        return "Failed to generate flash cards"


def call_deepseek_ai_quizes(summary_content, difficulty_level):
    Prompt = f"""
    You are an AI that generates quizzes based on a given summary. The quiz should be in JSON format with the following structure:

```json
    {{
        "quiz": {{
            "difficulty": "<easy | medium | hard>",
            "questions": [
            {{
                "question_text": "<A well-structured quiz question>",
                "choices": [
                "<Choice 1>",
                "<Choice 2>",
                "<Choice 3>",
                "<Choice 4>"
                ],
                "correct_answer": "<The correct choice from above>"
            }}
            ]
        }}
    }}
    ```

    Instructions:
    1. Extract key information** from the provided summary.
    2. Create at lest 3 and up to 20 based on the  given summary ltiple-choice questions based on that information.
    3. Each question should have 4 choices.
    4. The `"correct_answer"` key should contain the correct choice.
    5. Ensure the difficulty matches the requested level.

    Difficulty Level:
    - {difficulty_level}  

    ### **Expected JSON Output:**
    - The output should follow the quiz structure provided above, with `"difficulty"`, `"questions"`, and `"choices"`.
    - Make sure each question has exactly 4 answer choices and a correct answer.
    """
    try:
        # Initialize the OpenAI client with the DeepSeek API key and base URL
        client = OpenAI(api_key=settings.DEEPSEEK_API_KEY, base_url="https://api.deepseek.com")

        # Call the DeepSeek API to generate a summary
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": Prompt},
                {"role": "user", "content": summary_content},
            ],
            stream=False
        )
        # Extract the summary from the response
        logger.info(response.choices[0].message.content)
        quiz = response.choices[0].message.content
        return quiz

    except Exception as e:
        # Log the error and return a failure message
        logger.error("Failed to generate flash cards: {}".format(e))
        return "Failed to generate flash cards"


def clean_json_string(json_str):
    """
    Removes markdown-style JSON code block indicators (```json ... ```) if present.
    
    Args:
        json_str (str): The raw JSON string (possibly containing markdown).
    
    Returns:
        str: The cleaned JSON string.
    """
    # Remove markdown-style JSON code block
    cleaned_str = re.sub(r"```json|```", "", json_str).strip()
    
    return cleaned_str


def check_file_type(file_name):
    """
    Checks the file type based on the file name extension.
    
    Args:
        file_name (str): The name of the file.
    
    Returns:
        str: The file type based on the extension (e.g., 'pdf', 'docx', 'txt').
    """
    # Extract the file extension
    file_extension = file_name.split('.')[-1].lower()
    
    # Check the file type based on the extension
    if file_extension == 'pdf':
        return 'pdf'
    elif file_extension == 'docx':
        return 'docx'
    elif file_extension == 'txt':
        return 'txt'
    else:
        return 'unsupported'



def extract_pdf_text(file_path):
    """
    Extracts text content from a PDF file.
    Args:
        file_path (str): The file path to the PDF document.
    Returns:
        str: The extracted text from the PDF, with pages joined by newline characters.
             A separator line is appended at the end of the text.
    """
    try:
        text = []
        pdf = PdfReader(file_path)  # Read the PDF file content
        for page in pdf.pages:
            text.append(page.extract_text())
        # Join the text and append the separator
        return "\n".join(text) + "\n_________________________________\n"
    except Exception as e:
        # Log the error and return an empty string
        logger.error("Failed to extract text from PDF file: {}".format(e))
        return ""


def extract_docx_text(file_path):
    """
    Extracts text content from a DOCX file.
    Args:
        file_path (str): The file path to the DOCX document.
    Returns:
        str: The extracted text from the DOCX file.
    """
    try:
        # Read the DOCX file
        doc = Document(file_path)

        # Extract text from each paragraph
        text = [paragraph.text for paragraph in doc.paragraphs]

        # Join the text paragraphs
        return "\n".join(text) + "\n_________________________________\n"

    except Exception as e:
        # Log the error and return an empty string
        logger.error("Failed to extract text from DOCX file: {}".format(e))
        return ""


def extract_pptx_text(file_path):
    """
    Extracts text content from a PPTX file.
    Args:
        file_path (str): The file path to the PPTX document.
    Returns:
        str: The extracted text from the PPTX file.
    """
    try:
        # Read the PPTX file
        ppt = Presentation(file_path)

        # Extract text from each slide
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        # Join the text slides
        return "\n".join(text) + "\n_________________________________\n"

    except Exception as e:
        # Log the error and return an empty string
        logger.error("Failed to extract text from PPTX file: {}".format(e))
        return ""
